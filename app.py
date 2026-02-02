import sys
import os
import platform
from pathlib import Path
from datetime import datetime
import hashlib
import tempfile

import streamlit as st
import pandas as pd

# ----------------------------
# Ensure project root is importable (Streamlit Cloud safe)
# ----------------------------
ROOT = Path(__file__).resolve().parent
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

# ----------------------------
# Optional OCR dependencies (Cloud-safe)
# ----------------------------
OCR_AVAILABLE = True

try:
    import pytesseract  # type: ignore
except Exception:
    pytesseract = None
    OCR_AVAILABLE = False

try:
    from PIL import Image  # type: ignore
except Exception:
    Image = None

try:
    import pdf2image  # type: ignore
except Exception:
    pdf2image = None
    OCR_AVAILABLE = False

# ----------------------------
# Optional config
# ----------------------------
try:
    from config import TESSERACT_CMD, POPPLER_PATH, OCR_DPI, OCR_LANGUAGE
except Exception:
    TESSERACT_CMD = ""
    POPPLER_PATH = None
    OCR_DPI = 300
    OCR_LANGUAGE = "eng"

if pytesseract and TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

# ----------------------------
# App config
# ----------------------------
st.set_page_config(
    page_title="CUI Document Inspector",
    layout="wide"
)

st.title("📄 CUI Document Inspector")

# ============================================================================
# Sidebar (CLEAN — NO INDENTATION RISK)
# ============================================================================
st.sidebar.header("Navigation")
page = st.sidebar.radio(
    "Go to",
    [
        "Document Inspector",
        "Evidence Vault",
        "System Info",
    ]
)

st.sidebar.divider()
st.sidebar.markdown("### Runtime Checks")
st.sidebar.caption(f"OCR available: **{OCR_AVAILABLE}**")

st.sidebar.divider()
st.sidebar.markdown("### ℹ️ System Information")
st.sidebar.write(f"Python: {sys.version.split()[0]}")
st.sidebar.write(f"Platform: {platform.system()} {platform.release()}")

# ============================================================================
# Helper functions
# ============================================================================
def sha256_bytes(data: bytes) -> str:
    h = hashlib.sha256()
    h.update(data)
    return h.hexdigest()


def extract_text_from_pdf(uploaded_file) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(uploaded_file)
    text = ""

    for page in reader.pages:
        t = page.extract_text() or ""
        text += t

    # OCR fallback if scanned
    if not text.strip():
        if not OCR_AVAILABLE:
            st.warning(
                "⚠️ PDF appears scanned, but OCR is not available in this deployment."
            )
            return ""

        try:
            images = pdf2image.convert_from_bytes(
                uploaded_file.getvalue(),
                dpi=OCR_DPI,
                poppler_path=POPPLER_PATH,
            )
            for img in images:
                text += pytesseract.image_to_string(
                    img, lang=OCR_LANGUAGE
                )
        except Exception as e:
            st.error(f"OCR failed: {e}")

    return text


def extract_text_from_file(uploaded_file) -> str:
    name = uploaded_file.name.lower()

    if name.endswith(".pdf"):
        return extract_text_from_pdf(uploaded_file)

    if name.endswith(".txt"):
        return uploaded_file.getvalue().decode("utf-8", errors="ignore")

    if name.endswith(".docx"):
        from docx import Document
        doc = Document(uploaded_file)
        return "\n".join(p.text for p in doc.paragraphs)

    if name.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    return ""


# ============================================================================
# Pages
# ============================================================================
if page == "Document Inspector":
    st.header("Document Inspection")

    uploaded = st.file_uploader(
        "Upload a document",
        type=["pdf", "txt", "docx", "pptx"],
    )

    if uploaded:
        text = extract_text_from_file(uploaded)
        sha = sha256_bytes(uploaded.getvalue())

        st.subheader("File Metadata")
        st.write(
            {
                "filename": uploaded.name,
                "size_bytes": uploaded.size,
                "sha256": sha,
                "uploaded_at": datetime.utcnow().isoformat() + "Z",
            }
        )

        st.subheader("Extracted Text (preview)")
        st.text_area(
            "Preview",
            text[:5000],
            height=300,
        )

elif page == "Evidence Vault":
    st.header("Evidence Vault (Baseline)")
    st.info(
        "Evidence persistence will be added after baseline stabilization."
    )

elif page == "System Info":
    st.header("System Information")
    st.write("Python:", sys.version)
    st.write("Platform:", platform.platform())
    st.write("OCR available:", OCR_AVAILABLE)

